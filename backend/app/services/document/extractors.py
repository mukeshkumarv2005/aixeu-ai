"""Document text extractors — one per supported document type.

Each extractor implements ``ExtractorBase`` and is registered in the
``get_extractor()`` factory.
"""

from __future__ import annotations

import csv
import io
import os
import re
import subprocess
import tempfile
import traceback
from abc import ABC, abstractmethod
from pathlib import Path
from typing import IO, Any

from app.core.exceptions import AppException


# ── Exceptions ────────────────────────────────────────────────────────────────


class ExtractionError(AppException):
    """Raised when document text extraction fails irrecoverably."""

    def __init__(self, message: str, detail: dict | None = None) -> None:
        super().__init__(status_code=422, detail=message)
        self.detail_dict = detail or {}


class UnsupportedFormatError(AppException):
    """Raised when no extractor supports the given MIME type."""

    def __init__(self, mime_type: str) -> None:
        super().__init__(
            status_code=415,
            detail=f"Unsupported document format: {mime_type}",
        )


# ── Extractor result ──────────────────────────────────────────────────────────


class ExtractResult:
    """Result of a text extraction operation."""

    def __init__(
        self,
        text: str,
        page_count: int | None = None,
        ocr_used: bool = False,
        error_message: str | None = None,
    ) -> None:
        self.text = text
        self.page_count = page_count
        self.ocr_used = ocr_used
        self.error_message = error_message


# ── Base class ────────────────────────────────────────────────────────────────


class ExtractorBase(ABC):
    """Abstract base for all document extractors."""

    @abstractmethod
    def supported_mime_types(self) -> list[str]:
        """Return the MIME types this extractor can handle."""
        ...

    @abstractmethod
    async def extract(self, file_path: str, mime_type: str) -> ExtractResult:
        """Extract text content from the given file.

        Args:
            file_path: Absolute path to the file on disk.
            mime_type: MIME type of the file.

        Returns:
            An ``ExtractResult`` with the extracted text and metadata.

        Raises:
            ExtractionError: If extraction fails.
        """
        ...


# ── Text / Markdown extractor ─────────────────────────────────────────────────


class TextExtractor(ExtractorBase):
    """Extractor for plain text files (.txt) and Markdown (.md)."""

    def supported_mime_types(self) -> list[str]:
        return ["text/plain", "text/markdown"]

    async def extract(self, file_path: str, mime_type: str) -> ExtractResult:
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            return ExtractResult(text=text)
        except OSError as exc:
            raise ExtractionError(f"Cannot read text file: {exc}")


# ── CSV extractor ─────────────────────────────────────────────────────────────


class CsvExtractor(ExtractorBase):
    """Extractor for CSV files."""

    def supported_mime_types(self) -> list[str]:
        return ["text/csv", "application/csv"]

    async def extract(self, file_path: str, mime_type: str) -> ExtractResult:
        try:
            rows: list[str] = []
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append(" | ".join(row))
            text = "\n".join(rows)
            return ExtractResult(text=text)
        except OSError as exc:
            raise ExtractionError(f"Cannot read CSV file: {exc}")


# ── PDF extractor (pdfplumber primary, PyMuPDF fallback) ──────────────────────


class PdfExtractor(ExtractorBase):
    """Extractor for PDF documents.

    Uses ``pdfplumber`` as the primary extraction engine. Falls back to
    ``PyMuPDF`` (``fitz``) when pdfplumber produces empty or very short
    output, which can happen with certain scanned PDFs.
    """

    def supported_mime_types(self) -> list[str]:
        return ["application/pdf"]

    async def extract(self, file_path: str, mime_type: str) -> ExtractResult:
        text = ""
        page_count = 0
        errors: list[str] = []

        # ── Primary: pdfplumber ─────────────────────────────────────────
        try:
            import pdfplumber

            with pdfplumber.open(file_path) as pdf:
                pages = pdf.pages
                page_count = len(pages)
                page_texts: list[str] = []
                for page in pages:
                    page_text = page.extract_text() or ""
                    page_texts.append(page_text)
                text = "\n\n".join(page_texts)
        except ImportError:
            errors.append("pdfplumber not available")
        except Exception as exc:
            errors.append(f"pdfplumber failed: {exc}")

        # Fall back to PyMuPDF if pdfplumber produced nothing or failed
        if not text.strip():
            try:
                import fitz  # PyMuPDF

                doc = fitz.open(file_path)
                page_count = doc.page_count
                page_texts = []
                for page_num in range(doc.page_count):
                    page = doc.load_page(page_num)
                    page_texts.append(page.get_text())
                doc.close()
                text = "\n\n".join(page_texts)
            except ImportError:
                errors.append("PyMuPDF not available")
            except Exception as exc:
                errors.append(f"PyMuPDF failed: {exc}")

        error_msg = "; ".join(errors) if errors else None
        return ExtractResult(
            text=text,
            page_count=page_count,
            error_message=error_msg if not text.strip() else None,
        )


# ── DOCX extractor ────────────────────────────────────────────────────────────


class DocxExtractor(ExtractorBase):
    """Extractor for Word documents (.docx)."""

    def supported_mime_types(self) -> list[str]:
        return [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ]

    async def extract(self, file_path: str, mime_type: str) -> ExtractResult:
        try:
            from docx import Document

            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs]
            text = "\n".join(paragraphs)
            # Count approximate "pages" — each ~40 paragraphs ≈ 1 page
            page_count = max(1, len(paragraphs) // 40)
            return ExtractResult(text=text, page_count=page_count)
        except ImportError:
            raise ExtractionError("python-docx is not installed")
        except Exception as exc:
            raise ExtractionError(f"Cannot read DOCX file: {exc}")


# ── XLSX extractor ────────────────────────────────────────────────────────────


class XlsxExtractor(ExtractorBase):
    """Extractor for Excel spreadsheets (.xlsx)."""

    def supported_mime_types(self) -> list[str]:
        return [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ]

    async def extract(self, file_path: str, mime_type: str) -> ExtractResult:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True, data_only=True)
            parts: list[str] = []
            sheet_count = 0
            for sheet_name in wb.sheetnames:
                sheet_count += 1
                ws = wb[sheet_name]
                parts.append(f"--- Sheet: {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if row_text.strip():
                        parts.append(row_text)
            wb.close()
            text = "\n".join(parts)
            return ExtractResult(text=text, page_count=sheet_count)
        except ImportError:
            raise ExtractionError("openpyxl is not installed")
        except Exception as exc:
            raise ExtractionError(f"Cannot read XLSX file: {exc}")


# ── PPTX extractor ────────────────────────────────────────────────────────────


class PptxExtractor(ExtractorBase):
    """Extractor for PowerPoint presentations (.pptx)."""

    def supported_mime_types(self) -> list[str]:
        return [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ]

    async def extract(self, file_path: str, mime_type: str) -> ExtractResult:
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            parts: list[str] = []
            slide_count = 0
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_count += 1
                parts.append(f"--- Slide {slide_num} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            text = "\n".join(parts)
            return ExtractResult(text=text, page_count=slide_count)
        except ImportError:
            raise ExtractionError("python-pptx is not installed")
        except Exception as exc:
            raise ExtractionError(f"Cannot read PPTX file: {exc}")


# ── Image / OCR extractor ─────────────────────────────────────────────────────


class ImageExtractor(ExtractorBase):
    """Extractor for image files using OCR.

    Uses ``Pillow`` for image handling. OCR is performed via an external
    system dependency that can be swapped; by default, attempts to use
    ``tesserocr`` or the ``pytesseract`` wrapper.
    """

    # Common web-image MIME types
    MIME_TYPES = [
        "image/png",
        "image/jpeg",
        "image/webp",
        "image/gif",
        "image/bmp",
        "image/tiff",
    ]

    def supported_mime_types(self) -> list[str]:
        return self.MIME_TYPES

    async def extract(self, file_path: str, mime_type: str) -> ExtractResult:
        try:
            from PIL import Image as PILImage
        except ImportError:
            raise ExtractionError("Pillow is not installed")

        ocr_error: str | None = None
        text = ""

        # Attempt OCR with available engine
        text = self._ocr_with_pytesseract(file_path)
        if text.strip():
            return ExtractResult(text=text, ocr_used=True)

        text = self._ocr_with_tesserocr(file_path)
        if text.strip():
            return ExtractResult(text=text, ocr_used=True)

        # Fallback: try to extract EXIF text or just return empty
        ocr_error = (
            "No OCR engine available (install pytesseract or tesserocr)"
        )
        return ExtractResult(text="", ocr_used=False, error_message=ocr_error)

    def _ocr_with_pytesseract(self, file_path: str) -> str:
        """Attempt OCR via pytesseract."""
        try:
            import pytesseract
            from PIL import Image as PILImage

            img = PILImage.open(file_path)
            text = pytesseract.image_to_string(img)
            return text or ""
        except ImportError:
            return ""
        except Exception:
            return ""

    def _ocr_with_tesserocr(self, file_path: str) -> str:
        """Attempt OCR via tesserocr."""
        try:
            from PIL import Image as PILImage

            img = PILImage.open(file_path)
            # Try tesserocr directly
            import tesserocr

            api = tesserocr.PyTessBaseAPI()
            try:
                api.SetImage(img)
                text = api.GetUTF8Text()
                return text or ""
            finally:
                api.End()
        except ImportError:
            return ""
        except Exception:
            return ""


# ── Factory ───────────────────────────────────────────────────────────────────


_extractor_registry: dict[str, ExtractorBase] = {}


def _build_registry() -> dict[str, ExtractorBase]:
    """Build the MIME-type-to-extractor mapping."""
    extractors: list[ExtractorBase] = [
        TextExtractor(),
        CsvExtractor(),
        PdfExtractor(),
        DocxExtractor(),
        XlsxExtractor(),
        PptxExtractor(),
        ImageExtractor(),
    ]
    registry: dict[str, ExtractorBase] = {}
    for ext in extractors:
        for mt in ext.supported_mime_types():
            registry[mt] = ext
    return registry


def get_extractor(mime_type: str) -> ExtractorBase:
    """Return the appropriate extractor for the given MIME type.

    Raises:
        UnsupportedFormatError: If no extractor supports the MIME type.
    """
    global _extractor_registry
    if not _extractor_registry:
        _extractor_registry = _build_registry()

    extractor = _extractor_registry.get(mime_type)
    if extractor is None:
        # Try partial match for image types
        if mime_type.startswith("image/"):
            extractor = _extractor_registry.get("image/png")
    if extractor is None:
        raise UnsupportedFormatError(mime_type)
    return extractor
